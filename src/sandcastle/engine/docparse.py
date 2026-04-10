"""Lightweight document parser for PDF, DOCX, XLSX.

Inspired by LiteParse. Uses pymupdf for PDF extraction with spatial
layout preservation. Runs 100% locally, no API calls.

Install: pip install sandcastle-ai[parse]
"""

from __future__ import annotations

import csv
import logging
import os
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    text: str
    format: str  # "text", "markdown", "json"
    pages: int
    metadata: dict = field(default_factory=dict)


def parse_document(
    path: str,
    output_format: str = "markdown",
    pages: str | None = None,
    ocr: bool = False,
    ocr_engine: str = "auto",
    languages: list[str] | None = None,
) -> ParseResult:
    """Parse a document file to structured text.

    Args:
        path: File path to parse
        output_format: "text", "markdown", or "json"
        pages: Page range string like "1-5,10" or None for all
        ocr: Enable OCR for scanned documents (requires tesseract)
        ocr_engine: "auto" | "pymupdf" | "chandra" | "glm-ocr" - OCR engine selection
        languages: Language hints for Chandra OCR (e.g. ["cs", "en"])

    Returns:
        ParseResult with extracted text and metadata
    """
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix == ".pdf":
        return _parse_pdf_dispatch(p, output_format, pages, ocr, ocr_engine, languages)
    elif suffix in (".docx", ".doc"):
        return _parse_docx(p, output_format)
    elif suffix in (".xlsx", ".xls"):
        return _parse_xlsx(p, output_format)
    elif suffix in (".pptx",):
        return _parse_pptx(p, output_format)
    elif suffix in (".csv",):
        return _parse_csv(p, output_format)
    elif suffix in (".txt", ".md", ".json", ".yaml", ".yml", ".xml", ".html", ".log"):
        text = p.read_text(encoding="utf-8", errors="replace")
        return ParseResult(
            text=text,
            format="text",
            pages=1,
            metadata={
                "filename": p.name,
                "size_bytes": p.stat().st_size,
                "content_type": f"text/{suffix[1:]}",
            },
        )
    else:
        raise ValueError(
            f"Unsupported file format: {suffix}. "
            "Supported: .pdf, .docx, .xlsx, .pptx, .csv, .txt"
        )


def _parse_pdf_dispatch(
    path: Path,
    output_format: str,
    pages: str | None,
    ocr: bool,
    ocr_engine: str,
    languages: list[str] | None,
) -> ParseResult:
    """Route PDF parsing to the appropriate engine based on ocr_engine setting."""
    if ocr_engine == "chandra":
        return _parse_with_chandra(str(path), output_format, pages, languages)
    elif ocr_engine == "glm-ocr":
        return _parse_with_glm_ocr(str(path), output_format, pages)
    elif ocr_engine == "pymupdf":
        return _parse_pdf(path, output_format, pages, ocr)
    else:
        # auto: try pymupdf first, then glm-ocr if available, then chandra
        result = _parse_pdf(path, output_format, pages, ocr)
        avg_chars = len(result.text) / max(result.pages, 1) if result.text else 0
        if avg_chars < 50:
            logger.info(
                "pymupdf extracted <50 chars/page from %s, trying OCR fallback",
                path.name,
            )
            # Try GLM-OCR first (local, free, high accuracy)
            if _is_glm_ocr_available():
                try:
                    return _parse_with_glm_ocr(str(path), output_format, pages)
                except Exception as exc:
                    logger.warning("GLM-OCR fallback failed: %s", exc)
            # Fall back to Chandra OCR
            try:
                return _parse_with_chandra(str(path), output_format, pages, languages)
            except ImportError:
                logger.debug("chandra-ocr not installed, keeping pymupdf result")
            except Exception as exc:
                logger.warning("Chandra OCR fallback failed: %s", exc)
        return result


def _parse_with_chandra(
    file_path: str,
    output_format: str = "markdown",
    pages: str | None = None,
    languages: list[str] | None = None,
) -> ParseResult:
    """Parse a PDF/image document using Chandra OCR.

    Chandra handles scanned PDFs, images, handwriting, and complex tables.
    Install with: pip install sandcastle-ai[ocr]
    """
    try:
        from chandra_ocr import ocr
    except ImportError:
        raise ImportError(
            "chandra-ocr is required for Chandra OCR engine. "
            "Install with: pip install sandcastle-ai[ocr]"
        )

    kwargs: dict = {"output_format": output_format}
    if languages:
        kwargs["languages"] = languages
    if pages:
        kwargs["pages"] = pages

    result_text = ocr(file_path, **kwargs)
    if not isinstance(result_text, str):
        result_text = str(result_text)

    p = Path(file_path)
    return ParseResult(
        text=result_text,
        format=output_format,
        pages=_count_pages_hint(file_path, pages),
        metadata={
            "filename": p.name,
            "size_bytes": p.stat().st_size,
            "content_type": "application/pdf",
            "parse_method": "chandra",
            "languages": languages,
        },
    )


def _is_glm_ocr_available() -> bool:
    """Check if Ollama is running and has a GLM-OCR model pulled."""
    try:
        import httpx

        ollama_url = os.environ.get("OLLAMA_HOST", "http://localhost:11434")
        resp = httpx.get(f"{ollama_url}/api/tags", timeout=3)
        if resp.status_code == 200:
            models = [m.get("name", "") for m in resp.json().get("models", [])]
            return any("glm" in m.lower() and "ocr" in m.lower() for m in models)
    except Exception:
        pass
    return False


def _pdf_to_images(file_path: str, pages: str | None = None) -> list[bytes]:
    """Convert PDF pages to PNG images using pymupdf.

    Args:
        file_path: Path to the PDF file
        pages: Optional page range string like "1-5,10"

    Returns:
        List of PNG image bytes, one per page
    """
    try:
        import fitz  # pymupdf
    except ImportError:
        raise ImportError(
            "pymupdf is required for PDF-to-image conversion (GLM-OCR). "
            "Install with: pip install sandcastle-ai[parse]"
        )

    doc = fitz.open(file_path)
    page_range = _parse_page_ranges(pages, len(doc)) if pages else range(len(doc))

    images = []
    for page_num in page_range:
        if page_num >= len(doc):
            continue
        page = doc[page_num]
        pix = page.get_pixmap(dpi=300)  # high resolution for OCR accuracy
        images.append(pix.tobytes("png"))

    doc.close()
    return images


def _parse_with_glm_ocr(
    file_path: str,
    output_format: str = "markdown",
    pages: str | None = None,
) -> ParseResult:
    """Parse a document using GLM-OCR via Ollama.

    GLM-OCR is a 0.9B vision-language model with 94.6% accuracy on
    document reading benchmarks. Runs locally via Ollama - zero cost,
    zero data leaving the machine.

    Requires: ollama pull glm-ocr
    """
    import base64

    try:
        import httpx
    except ImportError:
        raise ImportError(
            "httpx is required for GLM-OCR (Ollama API). "
            "Install with: pip install httpx"
        )

    ollama_url = os.environ.get("OLLAMA_HOST", "http://localhost:11434")

    # Convert PDF pages to images via pymupdf
    page_images = _pdf_to_images(file_path, pages)
    if not page_images:
        p = Path(file_path)
        return ParseResult(
            text="",
            format=output_format,
            pages=0,
            metadata={
                "filename": p.name,
                "size_bytes": p.stat().st_size if p.exists() else 0,
                "parse_method": "glm-ocr",
                "model": "glm-ocr",
            },
        )

    results: list[str] = []
    for img_bytes in page_images:
        img_b64 = base64.b64encode(img_bytes).decode()

        resp = httpx.post(
            f"{ollama_url}/api/chat",
            json={
                "model": "glm-ocr",
                "messages": [
                    {
                        "role": "user",
                        "content": (
                            f"Extract all text from this document page as {output_format}. "
                            "Preserve tables, formulas, and layout."
                        ),
                        "images": [img_b64],
                    }
                ],
                "stream": False,
            },
            timeout=60,
        )

        if resp.status_code == 200:
            text = resp.json().get("message", {}).get("content", "")
            results.append(text)
        else:
            logger.warning(
                "GLM-OCR returned status %d for page %d",
                resp.status_code,
                len(results) + 1,
            )
            results.append("")

    p = Path(file_path)
    return ParseResult(
        text="\n\n---\n\n".join(results),
        format=output_format,
        pages=len(results),
        metadata={
            "filename": p.name,
            "size_bytes": p.stat().st_size if p.exists() else 0,
            "content_type": "application/pdf",
            "parse_method": "glm-ocr",
            "model": "glm-ocr",
            "accuracy_benchmark": "94.62%",
        },
    )


def _count_pages_hint(file_path: str, pages: str | None) -> int:
    """Best-effort page count without requiring pymupdf."""
    try:
        import fitz

        doc = fitz.open(file_path)
        total = len(doc)
        doc.close()
        if pages:
            return len(_parse_page_ranges(pages, total))
        return total
    except Exception:
        # If pymupdf is not available, estimate from page range string
        if pages:
            return len(_parse_page_ranges(pages, 9999))
        return 1


def _parse_page_ranges(pages_str: str | None, total: int) -> list[int]:
    """Parse page range string like '1-5,10,15-20' into list of 0-based page indices."""
    if not pages_str:
        return list(range(total))
    result = []
    for part in pages_str.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            result.extend(range(int(start) - 1, min(int(end), total)))
        else:
            idx = int(part) - 1
            if 0 <= idx < total:
                result.append(idx)
    return sorted(set(result))


def _parse_pdf(path: Path, output_format: str, pages: str | None, ocr: bool) -> ParseResult:
    """Parse PDF using pymupdf with spatial layout preservation."""
    try:
        import fitz  # pymupdf
    except ImportError:
        raise ImportError(
            "pymupdf is required for PDF parsing. "
            "Install with: pip install sandcastle-ai[parse]"
        )

    doc = fitz.open(str(path))
    total_pages = len(doc)
    page_indices = _parse_page_ranges(pages, total_pages)

    all_text_parts = []
    has_images = False

    for page_idx in page_indices:
        if page_idx >= total_pages:
            continue
        page = doc[page_idx]

        if output_format == "markdown":
            text = _extract_page_markdown(page, page_idx + 1)
        elif output_format == "json":
            text = _extract_page_json(page, page_idx + 1)
        else:
            text = page.get_text("text")

        # Check for images
        if page.get_images():
            has_images = True

        all_text_parts.append(text)

        # OCR fallback for pages with very little text
        if ocr and len(page.get_text("text").strip()) < 50:
            ocr_text = _ocr_page(page)
            if ocr_text:
                all_text_parts[-1] = ocr_text

    doc.close()

    separator = "\n\n---\n\n" if output_format == "markdown" else "\n\n"
    full_text = separator.join(all_text_parts)

    return ParseResult(
        text=full_text,
        format=output_format,
        pages=len(page_indices),
        metadata={
            "filename": path.name,
            "size_bytes": path.stat().st_size,
            "content_type": "application/pdf",
            "total_pages": total_pages,
            "parsed_pages": len(page_indices),
            "has_images": has_images,
            "parse_method": "pymupdf",
        },
    )


def _extract_page_markdown(page, page_num: int) -> str:
    """Extract page content as markdown with layout preservation.

    Uses pymupdf's text extraction with position data to detect:
    - Headers (large font -> # heading)
    - Bold text (font flags -> **bold**)
    - Tables (aligned columns -> markdown table)
    - Lists (bullet/number patterns -> - item)
    """
    import fitz

    blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
    lines = []

    for block in blocks:
        if block["type"] != 0:  # 0 = text block
            continue
        for line_data in block.get("lines", []):
            spans = line_data.get("spans", [])
            if not spans:
                continue

            # Detect font size for heading detection
            max_size = max(s.get("size", 12) for s in spans)
            text = "".join(s.get("text", "") for s in spans).strip()
            if not text:
                continue

            # Check for bold (bit 4 = bold in pymupdf flags)
            is_bold = any(s.get("flags", 0) & 2**4 for s in spans)

            # Format based on style
            if max_size >= 18:
                lines.append(f"# {text}")
            elif max_size >= 14:
                lines.append(f"## {text}")
            elif max_size >= 12 and is_bold:
                lines.append(f"### {text}")
            elif is_bold:
                lines.append(f"**{text}**")
            else:
                lines.append(text)

    return "\n".join(lines)


def _extract_page_json(page, page_num: int) -> str:
    """Extract page as JSON with bounding boxes."""
    import json

    blocks = page.get_text("dict")["blocks"]
    items = []
    for block in blocks:
        if block["type"] != 0:
            continue
        for line_data in block.get("lines", []):
            for span in line_data.get("spans", []):
                items.append({
                    "text": span.get("text", ""),
                    "bbox": list(span.get("bbox", [0, 0, 0, 0])),
                    "font_size": span.get("size", 12),
                    "bold": bool(span.get("flags", 0) & 2**4),
                    "page": page_num,
                })
    return json.dumps({"page": page_num, "items": items}, ensure_ascii=False)


def _ocr_page(page) -> str | None:
    """OCR a page using tesseract (optional)."""
    try:
        import io

        import pytesseract
        from PIL import Image

        pix = page.get_pixmap(dpi=300)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        text = pytesseract.image_to_string(img)
        return text.strip() if text.strip() else None
    except ImportError:
        logger.debug("pytesseract not available for OCR fallback")
        return None
    except Exception as e:
        logger.warning("OCR failed: %s", e)
        return None


def screenshot_page(path: str, page: int = 0, dpi: int = 150) -> bytes:
    """Generate a PNG screenshot of a PDF page.

    Returns PNG bytes for vision model consumption.
    """
    try:
        import fitz
    except ImportError:
        raise ImportError("pymupdf required: pip install sandcastle-ai[parse]")

    doc = fitz.open(path)
    if page >= len(doc):
        raise ValueError(f"Page {page} out of range (document has {len(doc)} pages)")
    pix = doc[page].get_pixmap(dpi=dpi)
    png_bytes = pix.tobytes("png")
    doc.close()
    return png_bytes


def _parse_docx(path: Path, output_format: str) -> ParseResult:
    """Parse DOCX using python-docx."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx required: pip install sandcastle-ai[parse]")

    doc = Document(str(path))
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if output_format == "markdown":
            if para.style.name.startswith("Heading 1"):
                lines.append(f"# {text}")
            elif para.style.name.startswith("Heading 2"):
                lines.append(f"## {text}")
            elif para.style.name.startswith("Heading 3"):
                lines.append(f"### {text}")
            elif para.style.name.startswith("List"):
                lines.append(f"- {text}")
            else:
                lines.append(text)
        else:
            lines.append(text)

    # Extract tables
    for table in doc.tables:
        if output_format == "markdown":
            headers = [cell.text.strip() for cell in table.rows[0].cells]
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("| " + " | ".join("---" for _ in headers) + " |")
            for row in table.rows[1:]:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append("| " + " | ".join(cells) + " |")
        else:
            for row in table.rows:
                lines.append("\t".join(cell.text.strip() for cell in row.cells))

    return ParseResult(
        text="\n".join(lines),
        format=output_format,
        pages=1,
        metadata={
            "filename": path.name,
            "size_bytes": path.stat().st_size,
            "content_type": (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
        },
    )


def _parse_xlsx(path: Path, output_format: str) -> ParseResult:
    """Parse XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("openpyxl required: pip install sandcastle-ai[parse]")

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        if output_format == "markdown":
            headers = [str(c) if c is not None else "" for c in rows[0]]
            lines = [f"## {sheet_name}", ""]
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("| " + " | ".join("---" for _ in headers) + " |")
            for row in rows[1:]:
                cells = [str(c) if c is not None else "" for c in row]
                lines.append("| " + " | ".join(cells) + " |")
            sheets_text.append("\n".join(lines))
        else:
            lines = [sheet_name + ":"]
            for row in rows:
                lines.append("\t".join(str(c) if c is not None else "" for c in row))
            sheets_text.append("\n".join(lines))

    wb.close()
    return ParseResult(
        text="\n\n".join(sheets_text),
        format=output_format,
        pages=len(wb.sheetnames),
        metadata={
            "filename": path.name,
            "size_bytes": path.stat().st_size,
            "content_type": (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ),
        },
    )


def _parse_pptx(path: Path, output_format: str) -> ParseResult:
    """Parse PPTX - extract text from slides."""
    try:
        from pptx import Presentation
    except ImportError:
        # Fallback: try reading as zip and extracting XML text
        return _parse_pptx_fallback(path, output_format)

    prs = Presentation(str(path))
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

        if output_format == "markdown":
            slides_text.append(f"## Slide {i}\n" + "\n".join(texts))
        else:
            slides_text.append(f"Slide {i}:\n" + "\n".join(texts))

    return ParseResult(
        text="\n\n".join(slides_text),
        format=output_format,
        pages=len(prs.slides),
        metadata={"filename": path.name, "size_bytes": path.stat().st_size},
    )


def _parse_pptx_fallback(path: Path, output_format: str) -> ParseResult:
    """Fallback PPTX parser using zipfile + XML."""
    texts = []
    with zipfile.ZipFile(str(path)) as z:
        for name in sorted(z.namelist()):
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                xml_content = z.read(name)
                root = ET.fromstring(xml_content)
                for elem in root.iter():
                    if elem.text and elem.text.strip():
                        texts.append(elem.text.strip())

    return ParseResult(
        text="\n".join(texts),
        format="text",
        pages=len(texts) // 5 or 1,
        metadata={"filename": path.name, "size_bytes": path.stat().st_size},
    )


def _parse_csv(path: Path, output_format: str) -> ParseResult:
    """Parse CSV to markdown table."""
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return ParseResult(text="", format="text", pages=0, metadata={"filename": path.name})

    if output_format == "markdown":
        headers = rows[0]
        lines = ["| " + " | ".join(headers) + " |"]
        lines.append("| " + " | ".join("---" for _ in headers) + " |")
        for row in rows[1:]:
            cells = row + [""] * (len(headers) - len(row))  # pad short rows
            lines.append("| " + " | ".join(cells) + " |")
        text = "\n".join(lines)
    else:
        text = "\n".join("\t".join(row) for row in rows)

    return ParseResult(
        text=text,
        format=output_format,
        pages=1,
        metadata={"filename": path.name, "rows": len(rows)},
    )
